"""Unit tests for PyMuPDFParserService (deterministic parsing, splitting helpers)."""
from __future__ import annotations

import io
import os


def _make_pdf_bytes(pages: list[str]) -> bytes:
    from reportlab.pdfgen import canvas

    buf = io.BytesIO()
    c = canvas.Canvas(buf)
    for i, text in enumerate(pages):
        c.drawString(72, 720, text)
        if i < len(pages) - 1:
            c.showPage()
    c.save()
    return buf.getvalue()


def test_parse_pdf_creates_page_chapters_and_markdown():
    from src.services.pymupdf_parser_service import PyMuPDFParserService

    prev = os.environ.pop("INGEST_MAX_DOC_CHUNKS", None)

    pdf_bytes = _make_pdf_bytes(["Hello page 1", "Hello page 2"])
    svc = PyMuPDFParserService()

    try:
        res = svc.parse_bytes(pdf_bytes, "sample.pdf", "application/pdf")
    finally:
        if prev is not None:
            os.environ["INGEST_MAX_DOC_CHUNKS"] = prev

    assert len(res.chapters) == 2
    assert res.chapters[0].heading == "Page 1"
    assert res.chapters[1].heading == "Page 2"
    assert "## Page 1" in res.markdown
    assert "Hello page 1" in res.markdown


def test_split_chapters_into_parts_splits_deterministically():
    from src.services.pymupdf_parser_service import ChapterEssence, split_chapters_into_parts

    import tiktoken
    enc = tiktoken.get_encoding("cl100k_base")

    chapters = [
        ChapterEssence(
            heading="Page 1",
            level=1,
            start_line=0,
            essence_text="Page 1\n\nA",
            full_content="A " * 50,
            images=[],
        ),
        ChapterEssence(
            heading="Page 2",
            level=1,
            start_line=1,
            essence_text="Page 2\n\nB",
            full_content="B " * 50,
            images=[],
        ),
        ChapterEssence(
            heading="Page 3",
            level=1,
            start_line=2,
            essence_text="Page 3\n\nC",
            full_content="C " * 50,
            images=[],
        ),
    ]

    t1 = len(enc.encode(chapters[0].full_content))
    t2 = len(enc.encode(chapters[1].full_content))

    parts = split_chapters_into_parts(
        chapters,
        "sample.pdf",
        max_tokens_per_part=max(1, t1 + t2 - 1),
    )

    assert len(parts) >= 2
    assert parts[0]["part_number"] == 1
    assert parts[0]["total_parts"] == len(parts)
    assert parts[0]["part_name"].startswith("sample_")


def test_parse_docx_extracts_text():
    from src.services.pymupdf_parser_service import PyMuPDFParserService

    from docx import Document

    buf = io.BytesIO()
    doc = Document()
    doc.add_paragraph("Hello DOCX")
    doc.save(buf)
    docx_bytes = buf.getvalue()

    svc = PyMuPDFParserService()
    res = svc.parse_bytes(docx_bytes, "sample.docx", "application/vnd.openxmlformats-officedocument.wordprocessingml.document")

    assert res.chapters
    assert "Hello DOCX" in res.markdown
    assert res.estimated_tokens >= 1


def test_parse_xlsx_extracts_cells():
    from src.services.pymupdf_parser_service import PyMuPDFParserService

    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws["A1"] = "Hello"
    ws["B1"] = "XLSX"
    buf = io.BytesIO()
    wb.save(buf)
    xlsx_bytes = buf.getvalue()

    svc = PyMuPDFParserService()
    res = svc.parse_bytes(xlsx_bytes, "sample.xlsx", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")

    assert res.chapters
    assert any("Sheet:" in ch.heading for ch in res.chapters)
    assert "Hello" in res.markdown
    assert "XLSX" in res.markdown
    assert res.estimated_tokens >= 1


def test_parse_pptx_extracts_slide_text():
    from src.services.pymupdf_parser_service import PyMuPDFParserService

    from pptx import Presentation
    from pptx.util import Inches

    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[5])
    tx = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    tx.text_frame.text = "Hello PPTX"
    buf = io.BytesIO()
    pres.save(buf)
    pptx_bytes = buf.getvalue()

    svc = PyMuPDFParserService()
    res = svc.parse_bytes(pptx_bytes, "sample.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation")

    assert res.chapters
    assert any(ch.heading.startswith("Slide ") for ch in res.chapters)
    assert "Hello PPTX" in res.markdown
    assert res.estimated_tokens >= 1


def test_parse_html_is_text_like():
    from src.services.pymupdf_parser_service import PyMuPDFParserService

    html = "<html><body><h1>Hello</h1><p>World</p></body></html>"
    svc = PyMuPDFParserService()
    res = svc.parse_bytes(html.encode("utf-8"), "sample.html", "text/html")

    assert res.chapters
    assert "<h1>Hello" in res.markdown
    assert res.estimated_tokens >= 1
