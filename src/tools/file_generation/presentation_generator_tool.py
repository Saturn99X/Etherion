"""
Production-ready PowerPoint Presentation Generation Tool using python-pptx.

Supports:
- Professional presentations with multiple slides
- Various slide layouts (title, content, two-column, etc.)
- Images, charts, and tables
- Custom themes and styling
- Speaker notes
- Bullet points and formatting
- Business presentations and reports

Library: python-pptx
Output: PowerPoint files (.pptx) with proper formatting
"""

import os
from typing import Dict, List, Optional, Any, Tuple
from io import BytesIO
from datetime import datetime
import logging

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

from .base_file_generator import BaseFileGenerator

logger = logging.getLogger(__name__)


class PresentationGeneratorTool(BaseFileGenerator):
    """
    Professional PowerPoint presentation generation tool using python-pptx.

    Features:
    - Multi-slide presentations
    - Various slide layouts
    - Images and charts
    - Tables and bullet points
    - Custom styling
    - Production-ready output
    """

    # Standard slide layouts
    LAYOUT_TITLE = 0
    LAYOUT_TITLE_CONTENT = 1
    LAYOUT_SECTION_HEADER = 2
    LAYOUT_TWO_CONTENT = 3
    LAYOUT_COMPARISON = 4
    LAYOUT_TITLE_ONLY = 5
    LAYOUT_BLANK = 6
    LAYOUT_CONTENT_CAPTION = 7
    LAYOUT_PICTURE_CAPTION = 8

    def __init__(
        self,
        tenant_id: str,
        agent_id: str,
        job_id: str,
        user_id: Optional[str] = None,
        project_id: Optional[str] = None,
    ):
        """
        Initialize the PowerPoint generator.

        Args:
            tenant_id: Tenant ID for isolation
            agent_id: Agent creating the presentation
            job_id: Associated job ID
            user_id: Optional user ID
            project_id: GCP project ID
        """
        super().__init__(tenant_id, agent_id, job_id, user_id, project_id)

        # Default colors (can be customized)
        self.primary_color = RGBColor(54, 96, 146)  # Dark blue
        self.secondary_color = RGBColor(79, 129, 189)  # Medium blue
        self.accent_color = RGBColor(192, 80, 77)  # Red accent
        self.text_color = RGBColor(51, 51, 51)  # Dark gray

        logger.info(f"Initialized PresentationGeneratorTool for tenant={tenant_id}")

    async def generate_presentation(
        self,
        slides: List[Dict[str, Any]],
        title: str = "Presentation",
        filename: str = "presentation.pptx",
        description: str = "",
        tags: Optional[List[str]] = None,
        metadata: Optional[Dict[str, Any]] = None,
    ) -> Dict[str, Any]:
        """
        Generate a complete PowerPoint presentation.

        Args:
            slides: List of slide configurations
            title: Presentation title
            filename: Output filename
            description: Description for metadata
            tags: Optional tags
            metadata: Optional additional metadata

        Returns:
            Dictionary containing asset info
        """
        try:
            # Create presentation
            prs = Presentation()

            # Set slide width and height (16:9 widescreen)
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(5.625)

            # Process each slide
            for slide_config in slides:
                self._add_slide(prs, slide_config)

            # Save to buffer
            buffer = BytesIO()
            prs.save(buffer)
            pptx_bytes = buffer.getvalue()
            buffer.close()

            # Save asset
            asset_info = await self.save_asset(
                file_bytes=pptx_bytes,
                filename=filename,
                mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                description=description or f"PowerPoint presentation: {title}",
                tags=tags or ["powerpoint", "presentation", "generated"],
                metadata={
                    **(metadata or {}),
                    "title": title,
                    "slide_count": len(slides),
                },
            )

            logger.info(f"Successfully generated PowerPoint presentation: {filename}")

            return asset_info

        except Exception as e:
            logger.error(f"Error generating presentation: {e}", exc_info=True)
            raise

    def _add_slide(self, prs: Presentation, slide_config: Dict[str, Any]):
        """
        Add a slide to the presentation based on configuration.

        Args:
            prs: Presentation object
            slide_config: Slide configuration dictionary
        """
        slide_type = slide_config.get("type", "title_content")
        content = slide_config.get("content", {})

        if slide_type == "title":
            self._add_title_slide(prs, content)
        elif slide_type == "title_content":
            self._add_title_content_slide(prs, content)
        elif slide_type == "two_column":
            self._add_two_column_slide(prs, content)
        elif slide_type == "table":
            self._add_table_slide(prs, content)
        elif slide_type == "section_header":
            self._add_section_header_slide(prs, content)
        elif slide_type == "blank":
            self._add_blank_slide(prs, content)
        else:
            # Default to title_content
            self._add_title_content_slide(prs, content)

    def _add_title_slide(self, prs: Presentation, content: Dict[str, Any]):
        """Add a title slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[self.LAYOUT_TITLE])

        # Set title
        title = slide.shapes.title
        title.text = content.get("title", "")
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = self.primary_color

        # Set subtitle
        if "subtitle" in content:
            subtitle = slide.placeholders[1]
            subtitle.text = content.get("subtitle", "")
            subtitle.text_frame.paragraphs[0].font.size = Pt(20)
            subtitle.text_frame.paragraphs[0].font.color.rgb = self.text_color

    def _add_title_content_slide(self, prs: Presentation, content: Dict[str, Any]):
        """Add a title and content slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[self.LAYOUT_TITLE_CONTENT])

        # Set title
        title = slide.shapes.title
        title.text = content.get("title", "")
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = self.primary_color

        # Set content
        if "bullet_points" in content:
            body = slide.placeholders[1]
            text_frame = body.text_frame
            text_frame.clear()

            for idx, point in enumerate(content["bullet_points"]):
                if isinstance(point, dict):
                    # Support nested bullet points
                    self._add_bullet_point(
                        text_frame,
                        point.get("text", ""),
                        level=point.get("level", 0),
                    )
                else:
                    self._add_bullet_point(text_frame, str(point), level=0)

        elif "text" in content:
            body = slide.placeholders[1]
            body.text = content["text"]

    def _add_two_column_slide(self, prs: Presentation, content: Dict[str, Any]):
        """Add a two-column slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[self.LAYOUT_TWO_CONTENT])

        # Set title
        title = slide.shapes.title
        title.text = content.get("title", "")
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.bold = True

        # Left column
        if "left_content" in content:
            left_placeholder = slide.placeholders[1]
            left_text = content["left_content"]
            if isinstance(left_text, list):
                text_frame = left_placeholder.text_frame
                text_frame.clear()
                for point in left_text:
                    self._add_bullet_point(text_frame, str(point), level=0)
            else:
                left_placeholder.text = str(left_text)

        # Right column
        if "right_content" in content:
            right_placeholder = slide.placeholders[2]
            right_text = content["right_content"]
            if isinstance(right_text, list):
                text_frame = right_placeholder.text_frame
                text_frame.clear()
                for point in right_text:
                    self._add_bullet_point(text_frame, str(point), level=0)
            else:
                right_placeholder.text = str(right_text)

    def _add_table_slide(self, prs: Presentation, content: Dict[str, Any]):
        """Add a slide with a table."""
        slide = prs.slides.add_slide(prs.slide_layouts[self.LAYOUT_TITLE_ONLY])

        # Set title
        title = slide.shapes.title
        title.text = content.get("title", "")
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.bold = True

        # Add table
        table_data = content.get("table_data", [])
        if table_data:
            rows = len(table_data)
            cols = len(table_data[0]) if table_data else 0

            # Calculate table position and size
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(9)
            height = Inches(3.5)

            # Add table shape
            table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
            table = table_shape.table

            # Fill table
            for row_idx, row_data in enumerate(table_data):
                for col_idx, cell_value in enumerate(row_data):
                    cell = table.cell(row_idx, col_idx)
                    cell.text = str(cell_value)

                    # Style header row
                    if row_idx == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = self.primary_color
                        paragraph = cell.text_frame.paragraphs[0]
                        paragraph.font.bold = True
                        paragraph.font.size = Pt(11)
                        paragraph.font.color.rgb = RGBColor(255, 255, 255)
                    else:
                        paragraph = cell.text_frame.paragraphs[0]
                        paragraph.font.size = Pt(10)

    def _add_section_header_slide(self, prs: Presentation, content: Dict[str, Any]):
        """Add a section header slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[self.LAYOUT_SECTION_HEADER])

        # Set title
        title = slide.shapes.title
        title.text = content.get("title", "")
        title.text_frame.paragraphs[0].font.size = Pt(40)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = self.primary_color
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Set subtitle
        if "subtitle" in content and len(slide.placeholders) > 1:
            subtitle = slide.placeholders[1]
            subtitle.text = content["subtitle"]
            subtitle.text_frame.paragraphs[0].font.size = Pt(20)
            subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def _add_blank_slide(self, prs: Presentation, content: Dict[str, Any]):
        """Add a blank slide with custom content."""
        slide = prs.slides.add_slide(prs.slide_layouts[self.LAYOUT_BLANK])

        # Add text boxes if specified
        if "text_boxes" in content:
            for text_box_config in content["text_boxes"]:
                self._add_text_box(slide, text_box_config)

    def _add_text_box(self, slide, config: Dict[str, Any]):
        """Add a text box to a slide."""
        left = Inches(config.get("left", 1))
        top = Inches(config.get("top", 1))
        width = Inches(config.get("width", 4))
        height = Inches(config.get("height", 1))

        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.text = config.get("text", "")

        # Apply formatting
        paragraph = text_frame.paragraphs[0]
        if "font_size" in config:
            paragraph.font.size = Pt(config["font_size"])
        if "bold" in config:
            paragraph.font.bold = config["bold"]
        if "color" in config:
            color = config["color"]
            if isinstance(color, str):
                # Convert hex color to RGB
                color = color.lstrip("#")
                r, g, b = tuple(int(color[i : i + 2], 16) for i in (0, 2, 4))
                paragraph.font.color.rgb = RGBColor(r, g, b)

    def _add_bullet_point(self, text_frame, text: str, level: int = 0):
        """Add a bullet point to a text frame."""
        if not text_frame.text:
            # First paragraph
            paragraph = text_frame.paragraphs[0]
        else:
            # Add new paragraph
            paragraph = text_frame.add_paragraph()

        paragraph.text = text
        paragraph.level = level
        paragraph.font.size = Pt(18 - (level * 2))  # Decrease size for nested levels
        paragraph.font.color.rgb = self.text_color

    async def generate_business_presentation(
        self,
        title: str,
        sections: List[Dict[str, Any]],
        filename: str = "business_presentation.pptx",
        author: Optional[str] = None,
        description: str = "",
        tags: Optional[List[str]] = None,
    ) -> Dict[str, Any]:
        """
        Generate a professional business presentation.

        Args:
            title: Presentation title
            sections: List of sections with content
            filename: Output filename
            author: Presentation author
            description: Description for metadata
            tags: Optional tags

        Returns:
            Dictionary containing asset info
        """
        try:
            slides = []

            # Title slide
            slides.append(
                {
                    "type": "title",
                    "content": {
                        "title": title,
                        "subtitle": f"Created by {author or 'Etherion AI'}\n{datetime.utcnow().strftime('%B %d, %Y')}",
                    },
                }
            )

            # Section slides
            for section in sections:
                section_type = section.get("type", "content")
                section_title = section.get("title", "")
                section_content = section.get("content", {})

                if section_type == "section_header":
                    slides.append(
                        {
                            "type": "section_header",
                            "content": {"title": section_title},
                        }
                    )
                elif section_type == "bullet_points":
                    slides.append(
                        {
                            "type": "title_content",
                            "content": {
                                "title": section_title,
                                "bullet_points": section_content.get("points", []),
                            },
                        }
                    )
                elif section_type == "table":
                    slides.append(
                        {
                            "type": "table",
                            "content": {
                                "title": section_title,
                                "table_data": section_content.get("data", []),
                            },
                        }
                    )
                elif section_type == "comparison":
                    slides.append(
                        {
                            "type": "two_column",
                            "content": {
                                "title": section_title,
                                "left_content": section_content.get("left", []),
                                "right_content": section_content.get("right", []),
                            },
                        }
                    )

            return await self.generate_presentation(
                slides=slides,
                title=title,
                filename=filename,
                description=description or f"Business presentation: {title}",
                tags=tags or ["powerpoint", "business", "presentation"],
                metadata={"author": author, "section_count": len(sections)},
            )

        except Exception as e:
            logger.error(f"Error generating business presentation: {e}", exc_info=True)
            raise

    async def generate_report_presentation(
        self,
        title: str,
        executive_summary: str,
        data_sections: List[Dict[str, Any]],
        conclusions: List[str],
        filename: str = "report_presentation.pptx",
        description: str = "",
        tags: Optional[List[str]] = None,
    ) -> Dict[str, Any]:
        """
        Generate a data-driven report presentation.

        Args:
            title: Report title
            executive_summary: Executive summary text
            data_sections: List of data sections with tables/charts
            conclusions: List of conclusion points
            filename: Output filename
            description: Description for metadata
            tags: Optional tags

        Returns:
            Dictionary containing asset info
        """
        try:
            slides = []

            # Title slide
            slides.append(
                {
                    "type": "title",
                    "content": {
                        "title": title,
                        "subtitle": f"Report Generated: {datetime.utcnow().strftime('%B %d, %Y')}",
                    },
                }
            )

            # Executive summary
            slides.append(
                {
                    "type": "title_content",
                    "content": {
                        "title": "Executive Summary",
                        "text": executive_summary,
                    },
                }
            )

            # Data sections
            for section in data_sections:
                if "table_data" in section:
                    slides.append(
                        {
                            "type": "table",
                            "content": {
                                "title": section.get("title", "Data"),
                                "table_data": section["table_data"],
                            },
                        }
                    )
                else:
                    slides.append(
                        {
                            "type": "title_content",
                            "content": {
                                "title": section.get("title", ""),
                                "text": section.get("text", ""),
                            },
                        }
                    )

            # Conclusions
            slides.append(
                {
                    "type": "title_content",
                    "content": {
                        "title": "Conclusions",
                        "bullet_points": conclusions,
                    },
                }
            )

            return await self.generate_presentation(
                slides=slides,
                title=title,
                filename=filename,
                description=description or f"Report presentation: {title}",
                tags=tags or ["powerpoint", "report", "data-driven"],
                metadata={"data_sections": len(data_sections)},
            )

        except Exception as e:
            logger.error(f"Error generating report presentation: {e}", exc_info=True)
            raise
